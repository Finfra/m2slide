#!/usr/bin/env python3
"""PPT 텍스트 강조·출처 추출 후 markdown 후처리.

Issue234 — ppt2m2slide 학습 round 3:
* PPT 컬러 텍스트 (빨강·파랑 등) → `**bold**` 자동 변환
* 출처 텍스트박스 ("출처:" / "[공통]" / URL 포함) → `::: source ... :::` 슬롯

Usage:
    python3 lib/ppt-emphasis-extract.py <pptx_path> <md_dir> [--dry-run]

핵심 메커니즘:
    1. python-pptx + lxml XML walk — `<a:rPr><a:solidFill><a:srgbClr>` 또는
       `<a:schemeClr val="accentN">` 으로 컬러 run 검출 (font.color.rgb 미지원 케이스 대응)
    2. 인접한 same-color run 병합 (PPT 가 단어 중간에서 run 분할하는 케이스 대응)
    3. 컬러 → markdown 매핑: 모든 비기본 컬러 → `**bold**` (Issue234 v1)
    4. 출처 검출: TEXT_BOX 타입 + URL 또는 "출처:/Source:/[공통]" prefix
    5. md 후처리: 정확 매칭으로 wrap, 코드블록·이미 감싼 곳 skip
"""
import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.exit("python-pptx not installed. pip install python-pptx")

NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# 강조로 인정하지 않을 색상 (본문·배경 기본)
SKIP_COLORS = {'000000', 'FFFFFF', 'tx1', 'tx2', 'bg1', 'bg2', 'dk1', 'dk2', 'lt1', 'lt2'}

URL_RE = re.compile(r'https?://\S+')
SOURCE_PREFIX_RE = re.compile(r'^\s*(출처|Source|source|\[공통\]|cf\)|참고|참조)\s*[:：]?')


def get_run_color(rpr_el):
    """Returns color identifier or None.
    srgbClr → hex string, schemeClr → "scheme:name", else None."""
    if rpr_el is None:
        return None
    srgb = rpr_el.find('.//a:srgbClr', NS)
    if srgb is not None:
        return srgb.get('val', '').upper()
    scheme = rpr_el.find('.//a:schemeClr', NS)
    if scheme is not None:
        return f"scheme:{scheme.get('val')}"
    return None


def extract_emphasis_and_sources(pptx_path: Path):
    """Returns (emphasis_phrases: list[str], source_blocks: list[str])."""
    prs = Presentation(str(pptx_path))
    emphasis_set = set()
    sources = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            full_text = (tf.text or '').strip()
            if not full_text:
                continue

            # 출처 텍스트박스 검출
            is_text_box = (shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX)
            has_url = bool(URL_RE.search(full_text))
            has_source_prefix = bool(SOURCE_PREFIX_RE.match(full_text))
            if (is_text_box and (has_url or has_source_prefix)) or has_source_prefix:
                sources.append(full_text)
                continue

            # 컬러 run 검출 + 인접 병합
            tf_el = tf._txBody
            for para_el in tf_el.findall('a:p', NS):
                merged_phrase = []
                last_color = None
                for r in para_el.findall('a:r', NS):
                    txt_el = r.find('a:t', NS)
                    txt = (txt_el.text or '') if txt_el is not None else ''
                    if not txt:
                        continue
                    color = get_run_color(r.find('a:rPr', NS))
                    is_colored = (
                        color is not None
                        and color not in SKIP_COLORS
                        and not color.startswith('scheme:tx')
                        and not color.startswith('scheme:bg')
                        and not color.startswith('scheme:dk')
                        and not color.startswith('scheme:lt')
                    )
                    if is_colored and color == last_color:
                        merged_phrase.append(txt)
                    else:
                        if merged_phrase and last_color is not None:
                            phrase = ''.join(merged_phrase).strip()
                            if phrase and len(phrase) >= 2:
                                emphasis_set.add(phrase)
                        merged_phrase = [txt] if is_colored else []
                        last_color = color if is_colored else None
                # paragraph 끝 — 남은 phrase flush
                if merged_phrase and last_color is not None:
                    phrase = ''.join(merged_phrase).strip()
                    if phrase and len(phrase) >= 2:
                        emphasis_set.add(phrase)

    return sorted(emphasis_set, key=len, reverse=True), sources


def wrap_emphasis_in_text(text: str, target: str) -> tuple[str, int]:
    """target 정확 매칭을 **target** 으로 변환. 코드블록·이미 감싼 케이스 skip."""
    parts = re.split(r'(```[\s\S]*?```)', text)
    count = 0
    for i, part in enumerate(parts):
        if part.startswith('```'):
            continue
        sub_parts = re.split(r'(`[^`\n]+`)', part)
        for j, sub in enumerate(sub_parts):
            if sub.startswith('`'):
                continue
            esc = re.escape(target)
            pattern = re.compile(rf'(?<!\*\*)({esc})(?!\*\*)')
            sub_parts[j], n = pattern.subn(r'**\1**', sub)
            count += n
        parts[i] = ''.join(sub_parts)
    return ''.join(parts), count


def convert_source_blockquote(text: str) -> tuple[str, int]:
    """`> 출처: ...` blockquote 라인 → `::: source\\n출처: ...\\n:::` 슬롯."""
    lines = text.split('\n')
    out = []
    count = 0
    i = 0
    while i < len(lines):
        line = lines[i]
        m = re.match(r'^>\s*(출처|Source|source|\[공통\]|cf\)|참고|참조)\s*[:：]?\s*(.+)$', line)
        if m:
            prefix, body = m.group(1), m.group(2)
            out.extend([
                '',
                '::: source',
                f'{prefix}: {body}' if not prefix.startswith('[') else f'{prefix} {body}',
                ':::',
                '',
            ])
            count += 1
            i += 1
            continue
        out.append(line)
        i += 1
    return '\n'.join(out), count


def patch_markdown(md_dir: Path, emphasis: list, dry_run: bool = False):
    md_files = sorted(md_dir.rglob('*.md'))
    changes = {'emphasis_applied': 0, 'source_applied': 0, 'files_changed': 0, 'emphasis_unmatched': []}
    matched_emphasis = set()

    for md_path in md_files:
        original = md_path.read_text(encoding='utf-8')
        text = original

        for em_text in emphasis:
            new_text, applied = wrap_emphasis_in_text(text, em_text)
            if applied:
                text = new_text
                changes['emphasis_applied'] += applied
                matched_emphasis.add(em_text)

        text, src_applied = convert_source_blockquote(text)
        changes['source_applied'] += src_applied

        if text != original:
            changes['files_changed'] += 1
            if not dry_run:
                md_path.write_text(text, encoding='utf-8')
            print(f"  patched: {md_path.relative_to(md_dir.parent)}")

    changes['emphasis_unmatched'] = [e for e in emphasis if e not in matched_emphasis]
    return changes


def main():
    ap = argparse.ArgumentParser(description='PPT 색 강조 + 출처 텍스트박스 → markdown 후처리')
    ap.add_argument('pptx', type=Path)
    ap.add_argument('md_dir', type=Path)
    ap.add_argument('--dry-run', action='store_true')
    args = ap.parse_args()

    if not args.pptx.is_file():
        sys.exit(f"PPT 파일 없음: {args.pptx}")
    if not args.md_dir.is_dir():
        sys.exit(f"md_dir 디렉토리 아님: {args.md_dir}")

    print(f"[1/2] PPT 분석: {args.pptx.name}")
    emphasis, sources = extract_emphasis_and_sources(args.pptx)
    print(f"  emphasis 후보: {len(emphasis)}개")
    for e in emphasis:
        print(f"    - {e[:80]}")
    print(f"  source 후보: {len(sources)}개")
    for s in sources:
        print(f"    - {s[:100]}")

    print(f"\n[2/2] markdown 후처리: {args.md_dir}")
    changes = patch_markdown(args.md_dir, emphasis, dry_run=args.dry_run)
    print(f"\n결과:")
    print(f"  파일 변경: {changes['files_changed']}개")
    print(f"  emphasis 적용: {changes['emphasis_applied']}건")
    print(f"  source 변환: {changes['source_applied']}건")
    if changes['emphasis_unmatched']:
        print(f"  emphasis 미매칭 ({len(changes['emphasis_unmatched'])}개) — md 텍스트와 정확 일치 없음:")
        for e in changes['emphasis_unmatched']:
            print(f"    × {e[:80]}")
    if args.dry_run:
        print("(dry-run — 파일 미저장)")


if __name__ == '__main__':
    main()
