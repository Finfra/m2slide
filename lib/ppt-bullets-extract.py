#!/usr/bin/env python3
"""PPT 슬라이드별 bullet paragraph 메타데이터 추출 — pptx2md bullet 손실 보완.

용도:
  pptx2md는 PPT 텍스트박스의 bullet markers(•/-/숫자)를 제거하고 plain paragraph로
  emit하는 케이스가 빈번 (특히 PowerPoint default bullet 사용 시). 결과적으로
  m2slide markdown 파서는 `<p>` 로 렌더하여 bullet 위계가 사라짐.

  본 스크립트는 python-pptx로 슬라이드별 모든 텍스트 paragraph의 bullet status를
  추출하여 manifest를 생성. ppt2m2slide agent가 이 정보를 참조하여 .md emit 시
  bullet markers(`*` / `  -`) 를 복원함.

bullet 판정 (python-pptx XML 직접 파싱):
  - <a:buChar char="•"/> 또는 <a:buAutoNum/> → bullet
  - <a:buNone/> → 명시적 no bullet
  - 미지정 → PowerPoint level default (대부분 bullet 활성)

출력:
  {out_dir}/bullets-manifest.yml — 슬라이드별·shape별·paragraph별 메타:
    slides:
      - n: 19
        shapes:
          - shape_idx: 1
            paragraphs:
              - { text: "사용법이 UNIX와 유사", level: 0, is_bullet: true, marker: "•" }
              - { text: "웹서비스 구현시 필요", level: 0, is_bullet: true, marker: "•" }

사용:
  python3 lib/ppt-bullets-extract.py <pptx_path> <out_dir>
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.stderr.write("ERROR: python-pptx not installed. Run: pip install python-pptx\n")
    sys.exit(1)

try:
    import yaml
except ImportError:
    sys.stderr.write("ERROR: PyYAML not installed. Run: pip install pyyaml\n")
    sys.exit(1)

# XML namespace
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def detect_bullet(para) -> dict:
    """Return {"is_bullet": bool, "marker": str|None, "explicit": bool}.

    Inspects <a:pPr> child elements for bullet definitions.
    """
    pPr = para._pPr
    if pPr is None:
        # No pPr → use PowerPoint default (most text boxes default to bullet)
        return {"is_bullet": True, "marker": None, "explicit": False}
    # Explicit <a:buNone/> → no bullet
    if pPr.find(A_NS + "buNone") is not None:
        return {"is_bullet": False, "marker": None, "explicit": True}
    # <a:buChar char="•"/> → bullet with char
    buChar = pPr.find(A_NS + "buChar")
    if buChar is not None:
        return {"is_bullet": True, "marker": buChar.get("char", "•"), "explicit": True}
    # <a:buAutoNum/> → numbered list
    buAutoNum = pPr.find(A_NS + "buAutoNum")
    if buAutoNum is not None:
        return {"is_bullet": True, "marker": "auto-num", "explicit": True}
    # No explicit bullet/no-bullet → inherit (assume bullet active)
    return {"is_bullet": True, "marker": None, "explicit": False}


def walk_shapes(shape, out):
    """Recursive walker for GROUP-nested text shapes."""
    if shape.has_text_frame:
        out.append(shape)
    elif shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            walk_shapes(sub, out)


def main():
    if len(sys.argv) != 3:
        sys.stderr.write(f"usage: {sys.argv[0]} <pptx_path> <out_dir>\n")
        sys.exit(2)

    pptx_path = Path(sys.argv[1])
    out_dir = Path(sys.argv[2])
    if not pptx_path.is_file():
        sys.stderr.write(f"ERROR: pptx not found: {pptx_path}\n")
        sys.exit(1)
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    manifest = {"source": str(pptx_path), "slides": []}
    total_paras = 0
    total_bullets = 0

    for n, slide in enumerate(prs.slides, 1):
        text_shapes = []
        for s in slide.shapes:
            walk_shapes(s, text_shapes)

        slide_entry = {"n": n, "shapes": []}
        for shape_idx, shape in enumerate(text_shapes, 1):
            tf = shape.text_frame
            paras_out = []
            for para in tf.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                bullet_info = detect_bullet(para)
                level = para.level or 0
                paras_out.append({
                    "text": text,
                    "level": level,
                    "is_bullet": bullet_info["is_bullet"],
                    "marker": bullet_info["marker"],
                    "bullet_explicit": bullet_info["explicit"],
                })
                total_paras += 1
                if bullet_info["is_bullet"]:
                    total_bullets += 1
            if paras_out:
                slide_entry["shapes"].append({
                    "shape_idx": shape_idx,
                    "left_emu": int(shape.left or 0),
                    "top_emu": int(shape.top or 0),
                    "width_emu": int(shape.width or 0),
                    "height_emu": int(shape.height or 0),
                    "paragraphs": paras_out,
                })
        manifest["slides"].append(slide_entry)

    manifest_path = out_dir / "bullets-manifest.yml"
    with open(manifest_path, "w", encoding="utf-8") as fh:
        yaml.safe_dump(manifest, fh, allow_unicode=True, sort_keys=False)

    sys.stderr.write(
        f"extracted {total_paras} paragraphs ({total_bullets} bullets) across {len(prs.slides)} slides\n"
    )
    sys.stderr.write(f"manifest: {manifest_path}\n")


if __name__ == "__main__":
    main()
