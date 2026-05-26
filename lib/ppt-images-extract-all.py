#!/usr/bin/env python3
"""PPT 슬라이드별 전체 이미지 추출 (pptx2md 보완 — 슬라이드당 1개 제한 우회).

용도:
  pptx2md는 슬라이드 1개당 대표 이미지 1개만 추출하므로, 다중 이미지 슬라이드
  (예: 3컬럼 예제 페이지)에서 이미지 손실 발생. 본 스크립트는 python-pptx로
  모든 PICTURE shape(중첩 GROUP 포함)를 추출하여 슬라이드+인덱스 네이밍으로 저장.

출력:
  - {out_dir}/s{slide_n:02d}_i{idx}.{ext}    — 이미지 파일 (slide 1-base, idx 1-base)
  - {out_dir}/manifest.yml                   — 슬라이드별 이미지 메타 (path, width, height, left, top)

사용:
  python3 lib/ppt-images-extract-all.py <pptx_path> <out_dir>
"""

import sys
import os
import hashlib
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    sys.stderr.write("ERROR: python-pptx not installed. Run: pip install python-pptx\n")
    sys.exit(1)

try:
    import yaml
except ImportError:
    sys.stderr.write("ERROR: PyYAML not installed. Run: pip install pyyaml\n")
    sys.exit(1)

try:
    from PIL import Image
    _HAS_PIL = True
except ImportError:
    _HAS_PIL = False

# Browser-incompatible formats → convert to PNG when possible
_RASTER_CONVERT = {"tif", "tiff", "bmp"}


def walk_pictures(shape, out):
    if shape.shape_type == 13:  # PICTURE
        out.append(shape)
    elif shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            walk_pictures(sub, out)


def main():
    if len(sys.argv) != 3:
        sys.stderr.write(f"usage: {sys.argv[0]} <pptx_path> <out_dir>\n")
        sys.exit(2)

    pptx_path = Path(sys.argv[1])
    out_dir = Path(sys.argv[2])
    if not pptx_path.is_file():
        sys.stderr.write(f"ERROR: pptx not found: {pptx_path}\n")
        sys.exit(1)
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    manifest = {"source": str(pptx_path), "slides": []}
    total_imgs = 0
    seen_blobs = {}  # sha1 → path (dedup)

    for n, slide in enumerate(prs.slides, 1):
        pics = []
        for s in slide.shapes:
            walk_pictures(s, pics)
        # sort top→bottom, left→right (reading order)
        pics.sort(key=lambda p: (p.top or 0, p.left or 0))

        slide_entry = {"n": n, "images": []}
        for idx, pic in enumerate(pics, 1):
            # GROUP-nested pictures fail pic.image accessor — fall back to part lookup
            try:
                blob = pic.image.blob
                ext = pic.image.ext
            except (AttributeError, KeyError):
                rId = pic._element.blip_rId
                image_part = slide.part.related_part(rId)
                blob = image_part.blob
                ct = image_part.content_type  # 'image/png', 'image/jpeg', ...
                ext = ct.split('/')[-1].replace('jpeg', 'jpg')
            sha1 = hashlib.sha1(blob).hexdigest()
            if sha1 in seen_blobs:
                # dedup — reuse existing file
                rel = seen_blobs[sha1]
            else:
                fname = f"s{n:02d}_i{idx}.{ext}"
                fpath = out_dir / fname
                with open(fpath, "wb") as fh:
                    fh.write(blob)
                # Convert browser-incompatible raster (tif/bmp) to PNG
                if ext.lower() in _RASTER_CONVERT and _HAS_PIL:
                    png_fname = f"s{n:02d}_i{idx}.png"
                    png_fpath = out_dir / png_fname
                    try:
                        Image.open(fpath).save(png_fpath, "PNG")
                        fpath.unlink()  # remove original .tif/.bmp
                        fname = png_fname
                    except Exception as e:
                        sys.stderr.write(f"WARN: failed to convert {fpath}: {e}\n")
                rel = fname
                seen_blobs[sha1] = rel
                total_imgs += 1
            slide_entry["images"].append({
                "idx": idx,
                "path": rel,
                "width_emu": int(pic.width or 0),
                "height_emu": int(pic.height or 0),
                "left_emu": int(pic.left or 0),
                "top_emu": int(pic.top or 0),
                "sha1": sha1,
            })
        manifest["slides"].append(slide_entry)

    manifest_path = out_dir / "manifest.yml"
    with open(manifest_path, "w", encoding="utf-8") as fh:
        yaml.safe_dump(manifest, fh, allow_unicode=True, sort_keys=False)

    sys.stderr.write(f"extracted {total_imgs} unique images across {len(prs.slides)} slides\n")
    sys.stderr.write(f"manifest: {manifest_path}\n")


if __name__ == "__main__":
    main()
